# app/domains/action/services/doc_summary_builder.py
"""
문서(PDF/DOCX/PPTX/HTML/MD/XLSX) 업로드 → 회의 정보 AI 추출 → PDF 생성 파이프라인.
"""
from __future__ import annotations

import io
import json
import logging
import re
import subprocess
import tempfile
import os
from datetime import datetime
from typing import Optional

from langchain_openai import ChatOpenAI
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import (
    HRFlowable,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.core.config import settings

logger = logging.getLogger(__name__)

# ── Korean CID fonts (built into reportlab, no system install needed) ─────────
pdfmetrics.registerFont(UnicodeCIDFont("HYGothic-Medium"))
pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))
_FONT = "HYGothic-Medium"
_FONT_BOLD = "HYSMyeongJo-Medium"

# ── Supported extensions → internal type ─────────────────────────────────────
EXT_MAP: dict[str, str] = {
    "pdf": "pdf",
    "pptx": "pptx",
    "ppt": "pptx",
    "html": "html",
    "htm": "html",
    "md": "md",
    "markdown": "md",
    "docx": "docx",
    "doc": "doc",
    "xlsx": "xlsx",
    "xls": "xls",
    "txt": "txt",
}

# ── LLM ──────────────────────────────────────────────────────────────────────
_llm = ChatOpenAI(model="gpt-4o-mini", api_key=settings.OPENAI_API_KEY)

_MAX_TEXT_CHARS = 14_000  # token budget ≈ 3500 tokens @ 4 chars/token

_EXTRACTION_PROMPT = """\
아래는 사용자가 업로드한 문서입니다. 문서 종류(회의록/기획서/보고서/메모 등)에 관계없이 \
회의 맥락으로 해석 가능한 내용을 추출해 JSON으로 정리하세요.

규칙:
- 불확실한 정보는 절대 추측하지 말고 "확인 필요"로 표기하세요.
- 내용이 없는 섹션은 빈 배열 []로 반환하세요.
- 미사여구 없이 실무용 요약만 작성하세요.
- 반드시 아래 JSON 스키마만 출력하세요(설명 텍스트 없이).

[문서]
{text}

[JSON 스키마]
{{
  "meeting_topic": "회의 주제 또는 추정 주제 | 확인 필요",
  "datetime": "회의 일시 | 확인 필요",
  "attendees": ["참석자 이름"],
  "discussion_items": [
    {{"topic": "논의 주제", "content": "논의 내용 요약"}}
  ],
  "decisions": [
    {{"content": "결정 사항"}}
  ],
  "action_items": [
    {{
      "assignee": "담당자 | 확인 필요",
      "content": "액션 내용",
      "deadline": "기한 | 확인 필요",
      "status": "미정"
    }}
  ],
  "issues_risks": [
    {{"content": "이슈 또는 리스크"}}
  ],
  "next_steps": ["다음 단계 내용"]
}}
"""


# ── Text extraction ───────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = [p.extract_text() for p in reader.pages if p.extract_text()]
    return "\n\n".join(p.strip() for p in pages if p.strip())


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_html(file_bytes: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _extract_md(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="replace")


def _extract_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="replace")


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _extract_doc_legacy(file_bytes: bytes) -> str:
    with tempfile.TemporaryDirectory() as tmpdir:
        doc_path = os.path.join(tmpdir, "input.doc")
        with open(doc_path, "wb") as f:
            f.write(file_bytes)
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, doc_path],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise ValueError(
                "구형 .doc 파일 변환에 실패했습니다. "
                "서버에 LibreOffice가 설치되어 있는지 확인하세요."
            )
        docx_path = os.path.join(tmpdir, "input.docx")
        with open(docx_path, "rb") as f:
            return _extract_docx(f.read())


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            sheets.append(f"[{sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _extract_xls_legacy(file_bytes: bytes) -> str:
    with tempfile.TemporaryDirectory() as tmpdir:
        xls_path = os.path.join(tmpdir, "input.xls")
        with open(xls_path, "wb") as f:
            f.write(file_bytes)
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", tmpdir, xls_path],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise ValueError(
                "구형 .xls 파일 변환에 실패했습니다. "
                "서버에 LibreOffice가 설치되어 있는지 확인하세요."
            )
        xlsx_path = os.path.join(tmpdir, "input.xlsx")
        with open(xlsx_path, "rb") as f:
            return _extract_xlsx(f.read())


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = (filename or "").rsplit(".", 1)[-1].lower() if "." in filename else ""
    file_type = EXT_MAP.get(ext)
    if not file_type:
        raise ValueError(
            f"지원하지 않는 파일 형식: .{ext}  "
            f"(지원 형식: {', '.join(sorted(EXT_MAP))})"
        )

    if file_type == "pdf":
        text = _extract_pdf(file_bytes)
    elif file_type == "pptx":
        text = _extract_pptx(file_bytes)
    elif file_type == "html":
        text = _extract_html(file_bytes)
    elif file_type == "md":
        text = _extract_md(file_bytes)
    elif file_type == "docx":
        text = _extract_docx(file_bytes)
    elif file_type == "doc":
        text = _extract_doc_legacy(file_bytes)
    elif file_type == "xlsx":
        text = _extract_xlsx(file_bytes)
    elif file_type == "xls":
        text = _extract_xls_legacy(file_bytes)
    elif file_type == "txt":
        text = _extract_txt(file_bytes)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {file_type}")

    if not text.strip():
        raise ValueError(
            "텍스트를 추출할 수 없습니다. "
            "스캔 이미지 PDF의 경우 OCR 처리가 필요하며 현재 지원하지 않습니다."
        )
    return text


# ── AI extraction ─────────────────────────────────────────────────────────────

async def extract_meeting_info(text: str) -> dict:
    truncated = text[:_MAX_TEXT_CHARS]
    prompt = _EXTRACTION_PROMPT.format(text=truncated)

    result = await _llm.ainvoke(prompt)
    raw = result.content

    json_match = re.search(r"\{.*\}", raw, re.DOTALL)
    if not json_match:
        raise ValueError(
            "AI 응답에서 구조화된 데이터를 파싱할 수 없습니다. "
            "잠시 후 다시 시도해 주세요."
        )
    try:
        return json.loads(json_match.group())
    except json.JSONDecodeError as exc:
        logger.warning("AI JSON parse error: %s | raw=%s", exc, raw[:200])
        raise ValueError(
            "AI 응답 형식이 올바르지 않습니다. 잠시 후 다시 시도해 주세요."
        )


# ── PDF generation ────────────────────────────────────────────────────────────

_PRIMARY = colors.HexColor("#5668F3")
_SECONDARY = colors.HexColor("#EEF0FE")
_BORDER = colors.HexColor("#D0D5F5")
_GRAY = colors.HexColor("#888899")


def _style(name: str, **kw) -> ParagraphStyle:
    defaults = dict(fontName=_FONT, fontSize=10, leading=16)
    defaults.update(kw)
    return ParagraphStyle(name, **defaults)


def generate_pdf(summary: dict, original_filename: str) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        rightMargin=2.2 * cm,
        leftMargin=2.2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    title_s = _style("Title", fontName=_FONT_BOLD, fontSize=20, textColor=_PRIMARY, spaceAfter=6)
    meta_s = _style("Meta", fontSize=10, textColor=_GRAY, spaceAfter=3)
    section_s = _style("Section", fontName=_FONT_BOLD, fontSize=13, textColor=_PRIMARY, spaceBefore=14, spaceAfter=8)
    body_s = _style("Body", fontSize=10, spaceAfter=4)
    bullet_s = _style("Bullet", fontSize=10, leftIndent=14, spaceAfter=3)
    sub_topic_s = _style("SubTopic", fontName=_FONT_BOLD, fontSize=10, spaceAfter=2)
    footer_s = _style("Footer", fontSize=8, textColor=_GRAY, spaceBefore=4)

    story = []

    # ── 헤더 ─────────────────────────────────────────────────────────────────
    story.append(Paragraph("회의 요약 보고서", title_s))

    meeting_topic = summary.get("meeting_topic") or "확인 필요"
    dt = summary.get("datetime") or "확인 필요"
    attendees: list = summary.get("attendees") or []

    story.append(Paragraph(f"주제: {meeting_topic}", meta_s))
    story.append(Paragraph(f"일시: {dt}", meta_s))
    if attendees:
        story.append(Paragraph(f"참석자: {', '.join(attendees)}", meta_s))
    story.append(Paragraph(f"원본 파일: {original_filename}", meta_s))
    story.append(HRFlowable(width="100%", color=_PRIMARY, thickness=2, spaceAfter=14))

    # ── 논의 사항 ─────────────────────────────────────────────────────────────
    discussion: list = summary.get("discussion_items") or []
    if discussion:
        story.append(Paragraph("논의 사항", section_s))
        for i, item in enumerate(discussion, 1):
            topic = item.get("topic") or ""
            content = item.get("content") or ""
            story.append(Paragraph(f"{i}. {topic}", sub_topic_s))
            if content:
                story.append(Paragraph(content, bullet_s))

    # ── 결정 사항 ─────────────────────────────────────────────────────────────
    decisions: list = summary.get("decisions") or []
    if decisions:
        story.append(Paragraph("결정 사항", section_s))
        for d in decisions:
            story.append(Paragraph(f"• {d.get('content', '')}", bullet_s))

    # ── 액션 아이템 (표) ──────────────────────────────────────────────────────
    action_items: list = summary.get("action_items") or []
    if action_items:
        story.append(Paragraph("액션 아이템", section_s))

        header = ["담당자", "내용", "기한", "상태"]
        rows = [header] + [
            [
                a.get("assignee") or "확인 필요",
                a.get("content") or "",
                a.get("deadline") or "확인 필요",
                a.get("status") or "미정",
            ]
            for a in action_items
        ]
        col_w = [3.2 * cm, 9.5 * cm, 3 * cm, 1.8 * cm]
        tbl = Table(rows, colWidths=col_w, repeatRows=1)
        tbl.setStyle(
            TableStyle([
                ("FONTNAME", (0, 0), (-1, 0), _FONT_BOLD),
                ("FONTNAME", (0, 1), (-1, -1), _FONT),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("BACKGROUND", (0, 0), (-1, 0), _PRIMARY),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, _SECONDARY]),
                ("GRID", (0, 0), (-1, -1), 0.5, _BORDER),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("WORDWRAP", (1, 1), (1, -1), True),
            ])
        )
        story.append(tbl)

    # ── 이슈 / 리스크 ─────────────────────────────────────────────────────────
    issues: list = summary.get("issues_risks") or []
    if issues:
        story.append(Paragraph("이슈 / 리스크", section_s))
        for iss in issues:
            story.append(Paragraph(f"• {iss.get('content', '')}", bullet_s))

    # ── 후속 일정 / 다음 단계 ─────────────────────────────────────────────────
    next_steps: list = summary.get("next_steps") or []
    if next_steps:
        story.append(Paragraph("후속 일정 / 다음 단계", section_s))
        for ns in next_steps:
            story.append(Paragraph(f"• {ns}", bullet_s))

    # ── 내용 없음 안내 ────────────────────────────────────────────────────────
    has_content = any([discussion, decisions, action_items, issues, next_steps])
    if not has_content:
        story.append(Spacer(1, 0.5 * cm))
        story.append(Paragraph(
            "문서에서 회의 관련 내용을 찾을 수 없습니다. "
            "원본 문서를 확인해 주세요.",
            _style("NoContent", fontSize=11, textColor=_GRAY),
        ))

    # ── 푸터 ─────────────────────────────────────────────────────────────────
    story.append(Spacer(1, 0.6 * cm))
    story.append(HRFlowable(width="100%", color=_BORDER, thickness=0.5))
    story.append(Paragraph(
        f"생성일시: {datetime.now().strftime('%Y-%m-%d %H:%M')} | AI 자동 추출 요약 — 중요 내용은 반드시 원본 문서에서 확인하세요.",
        footer_s,
    ))

    doc.build(story)
    buf.seek(0)
    return buf.read()


# ── Main pipeline ─────────────────────────────────────────────────────────────

async def process_document(file_bytes: bytes, filename: str) -> tuple[bytes, str]:
    """
    전체 파이프라인: 텍스트 추출 → AI 추출 → PDF 생성.
    Returns (pdf_bytes, output_filename).
    """
    text = extract_text(file_bytes, filename)
    summary = await extract_meeting_info(text)
    pdf_bytes = generate_pdf(summary, filename)

    ts = datetime.now().strftime("%Y%m%d_%H%M")
    output_filename = f"meeting_summary_{ts}.pdf"
    return pdf_bytes, output_filename
